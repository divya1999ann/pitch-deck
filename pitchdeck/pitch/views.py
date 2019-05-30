from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.models import User
from django.http import HttpResponse
from django.shortcuts import render
from django.shortcuts import render_to_response
from django.template import RequestContext
from pitch.models import userprofile, startup_data
from pptx import Presentation

from pitch.forms import UserForm, userprofileform, startup_Form


# Create your views here.
def blank(request):
    return render(request, 'pitch/blank.html')

def home(request):
    return render(request,'pitch/home.html')
def register(request):
    context = RequestContext(request)
    registered = False

    if request.method == 'POST':
        user_form = UserForm(data=request.POST)
        user_profile = userprofileform(data=request.POST)
        if user_form.is_valid() and user_profile.is_valid():
            request.session['username'] = "alpha"

            user = user_form.save()
            user.set_password(user.password)
            user.save()
            profile = user_profile.save(commit=False)
            profile.user = user
            profile.save()
            if 'picture' in request.FILES:
                profile.picture = request.FILES['picture']

                # Now we save the UserProfile model instance.
            profile.save()

            registered = True
        else:
            print(user_form.errors, user_profile.errors)
    else:
        user_form = UserForm
        user_profile = userprofileform

    return render(request, "pitch/register.html",
                  {'user_form': user_form, 'user_profile': user_profile, 'registered': registered, })


def user_login(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        user = authenticate(request, username=username, password=password)

        if user is not None:
            login(request, user)
            # request.session['fav_color'] = 'blue';
            # request.session.modified = True

            data = User.objects.get(username=username)
            info = userprofile.objects.get(user_id=data.id)
            request.session['username'] = username
            context_dict = {'profile': data, 'info': info}
            return render(request, 'pitch/loginhome.html', context_dict)
        else:
            print("Invalid login details :{0},{1}".format(username, password))
            return HttpResponse("Invalid login details supplied.")
    else:
        return render(request, 'pitch/login.html', {})


def input(request):
    if request.POST:
        form = startup_Form(request.POST)
        if form.is_valid():
            form.save()
    form = startup_Form()
    context = {
        'form': form
    }
    return render(request, 'pitch/input.html', context)

def create(request):
    if request.session.has_key('username'):
        username = request.session['username']

        create_pitchdeck(username)
    return HttpResponse("p")


def create_pitchdeck(username):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    data = User.objects.get(username=username)
    title1 = data.username
    sbtitle1 = data.email

    title.text = title1
    subtitle.text = sbtitle1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    print("Slide 2\n")

    title2 = data.email
    title_shape.text = title2

    tf = body_shape.text_frame
    sbtitle2 = data.email
    tf.text = sbtitle2
    point1 = data.email
    p = tf.add_paragraph()
    p.text = point1
    p.level = 1
    sbpoint1 = data.email
    p = tf.add_paragraph()
    p.text = sbpoint1
    p.level = 2
    prs.save(data.username + '.pptx')

def user_logout(request):
    logout(request)
    return HttpResponse('you have logged out')


def ppt_generate(i):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title1 = i.username
    sbtitle1 = i.email

    title.text = title1
    subtitle.text = sbtitle1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    print("Slide 2\n")

    title2 = i.username
    title_shape.text = title2

    tf = body_shape.text_frame
    sbtitle2 = i.email
    tf.text = sbtitle2
    point1 = i.username
    p = tf.add_paragraph()
    p.text = point1
    p.level = 1
    sbpoint1 = i.email
    p = tf.add_paragraph()
    p.text = sbpoint1
    p.level = 2

    prs.save(i.username + '.pptx')


